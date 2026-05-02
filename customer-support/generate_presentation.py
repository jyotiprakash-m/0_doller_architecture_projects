import os
import re
import base64
import requests
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Configuration & Colors ---
INDIGO = RGBColor(99, 102, 241)   # Frontend
EMERALD = RGBColor(16, 185, 129)  # Backend
AMBER = RGBColor(245, 158, 11)    # Datastore
PURPLE = RGBColor(168, 85, 247)   # AI/LangGraph
PINK = RGBColor(236, 72, 153)     # Worker/Kafka
SLATE = RGBColor(30, 41, 59)      # Text/Background

ARCH_FILE = "ARCHITECTURE.md"

def render_mermaid(mmd_code):
    """Renders Mermaid code to an image using mermaid.ink."""
    try:
        # Clean up code (remove leading/trailing whitespace and extra newlines)
        mmd_code = mmd_code.strip()
        # Mermaid.ink expects base64 encoded string
        # We use standard base64 encoding
        sample_string_bytes = mmd_code.encode("utf-8")
        base64_bytes = base64.b64encode(sample_string_bytes)
        base64_string = base64_bytes.decode("ascii")
        
        url = f"https://mermaid.ink/img/{base64_string}"
        response = requests.get(url)
        if response.status_code == 200:
            return io.BytesIO(response.content)
        else:
            print(f"Failed to render diagram. Status code: {response.status_code}")
            return None
    except Exception as e:
        print(f"Error rendering mermaid: {e}")
        return None

def extract_mermaid_blocks(file_path):
    """Extracts all mermaid code blocks from a markdown file."""
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return []
    
    with open(file_path, "r") as f:
        content = f.read()
    
    # Regex to find ```mermaid blocks
    blocks = re.findall(r"```mermaid\n(.*?)\n```", content, re.DOTALL)
    return blocks

def add_title_subtitle(slide, title_text, subtitle_text=""):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = INDIGO
    title.text_frame.paragraphs[0].font.bold = True

    if subtitle_text:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = SLATE

def create_presentation():
    prs = Presentation()
    mmd_blocks = extract_mermaid_blocks(ARCH_FILE)
    
    print(f"Found {len(mmd_blocks)} Mermaid diagrams in {ARCH_FILE}")

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SupportSim AI"
    subtitle.text = "The Future of Private AI Customer Support Training\nFull Architecture & System Design"
    title.text_frame.paragraphs[0].font.color.rgb = INDIGO
    title.text_frame.paragraphs[0].font.size = Pt(44)

    # --- Slide 2: Mission & Vision ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "System Mission & Vision")
    tf = slide.placeholders[1].text_frame
    tf.text = "A local-first, privacy-focused SaaS platform for training support agents."
    bullets = [
        "🛡️ 100% Data Privacy: Local AI inference via Ollama.",
        "💰 Zero API Costs: High-performance open LLMs (Llama 3.2).",
        "⚙️ Enterprise Scalability: Kafka-backed distributed processing.",
        "🤖 Intelligent Agents: Multi-agent state management with LangGraph."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0

    # --- Slide 3: Core Features ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "Core Platform Features")
    tf = slide.placeholders[1].text_frame
    features = [
        "Advanced RAG: Multi-tenant document isolation.",
        "Dynamic Simulation: Evolving emotional personas.",
        "Automated Evaluation: Multi-metric performance scoring.",
        "SaaS Monetization: Integrated Stripe billing & credits."
    ]
    for feat in features:
        p = tf.add_paragraph()
        p.text = f"• {feat}"

    # --- Slide 4: High-Level Architecture (Diagram) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "High-Level Architecture")
    if len(mmd_blocks) > 0:
        img_stream = render_mermaid(mmd_blocks[0])
        if img_stream:
            slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), height=Inches(5.5))
    else:
        slide.placeholders[1].text = "Diagram placeholder: High-Level System Flow"

    # --- Slide 5: Technical Stack ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "Modern Tech Stack")
    tf = slide.placeholders[1].text_frame
    stack = [
        "Frontend: Next.js 15, Tailwind CSS, Framer Motion",
        "Backend: FastAPI (Python 3.11+), JWT Auth",
        "AI Engine: Ollama, LangGraph, LlamaIndex",
        "Data: SQLite (Relational), ChromaDB (Vector)",
        "Messaging: Kafka Message Broker"
    ]
    for tech in stack:
        p = tf.add_paragraph()
        p.text = f"• {tech}"

    # --- Slide 6: Document Indexing Workflow (Diagram) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "Document Indexing Pipeline (Kafka)")
    if len(mmd_blocks) > 1:
        img_stream = render_mermaid(mmd_blocks[1])
        if img_stream:
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))
    else:
        slide.placeholders[1].text = "Diagram placeholder: Kafka Indexing Sequence"

    # --- Slide 7: Training Session Workflow (Diagram) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "Training Session Workflow (LangGraph)")
    if len(mmd_blocks) > 2:
        img_stream = render_mermaid(mmd_blocks[2])
        if img_stream:
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))

    # --- Slide 8: Billing & Credits (Stripe Diagram) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "Billing & Credit Fulfillment (Stripe)")
    if len(mmd_blocks) > 3:
        img_stream = render_mermaid(mmd_blocks[3])
        if img_stream:
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))

    # --- Slide 9: Database Schema (ERD Diagram) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "Relational Database Schema (ERD)")
    if len(mmd_blocks) > 4:
        img_stream = render_mermaid(mmd_blocks[4])
        if img_stream:
            # ERDs are usually tall, resize accordingly
            slide.shapes.add_picture(img_stream, Inches(2), Inches(1.5), height=Inches(5.5))

    # --- Slide 10: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_subtitle(slide, "Conclusion: The Private AI Edge")
    tf = slide.placeholders[1].text_frame
    tf.text = "SupportSim AI provides a premium, scalable foundation for enterprise training without compromising data privacy."
    p = tf.add_paragraph()
    p.text = "\nThank You!"
    p.alignment = PP_ALIGN.CENTER

    # Save Presentation
    output_path = "SupportSim_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"Success! Presentation saved to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    print("Generating presentation with Mermaid diagrams...")
    create_presentation()
